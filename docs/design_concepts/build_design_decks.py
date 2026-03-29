#!/usr/bin/env python3
"""Build PPTX files, each with one full-slide design concept image."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

HERE = Path(__file__).resolve().parent
W = Inches(13.333)
H = Inches(7.5)

DECKS = [
    (
        "EZ_Trainz_Design_01_Glassmorphism.pptx",
        HERE / "ez_trainz_design_glassmorphism.png",
        "Concept 1 — Glassmorphism & soft gradients",
    ),
    (
        "EZ_Trainz_Design_02_Playful.pptx",
        HERE / "ez_trainz_design_playful.png",
        "Concept 2 — Bold playful edtech",
    ),
    (
        "EZ_Trainz_Design_03_Minimal_Dark.pptx",
        HERE / "ez_trainz_design_minimal_dark.png",
        "Concept 3 — Minimal premium dark",
    ),
    (
        "EZ_Trainz_Design_04_Neumorphism.pptx",
        HERE / "ez_trainz_design_neumorphism.png",
        "Concept 4 — Neumorphism soft UI",
    ),
    (
        "EZ_Trainz_Design_05_Brutalist.pptx",
        HERE / "ez_trainz_design_brutalist.png",
        "Concept 5 — Brutalist editorial",
    ),
    (
        "EZ_Trainz_Design_06_Zen_Paper.pptx",
        HERE / "ez_trainz_design_zen_paper.png",
        "Concept 6 — Zen paper & calm study",
    ),
]


def main() -> None:
    for filename, img_path, caption in DECKS:
        if not img_path.is_file():
            raise SystemExit(f"Missing image: {img_path}")

        prs = Presentation()
        prs.slide_width = W
        prs.slide_height = H

        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(img_path), 0, 0, width=W, height=H)

        # Caption bar at bottom (readable on any background)
        box = slide.shapes.add_textbox(Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.55))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = caption
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 30, 30)

        out = HERE / filename
        prs.save(str(out))
        print(f"Wrote {out}")


if __name__ == "__main__":
    main()
